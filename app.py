from flask import Flask, request, jsonify, render_template
from flask_cors import CORS
import fitz  # PyMuPDF
import os
import requests
from dotenv import load_dotenv
from sentence_transformers import SentenceTransformer, util
from pptx import Presentation
import docx
import textract
import re

app = Flask(__name__, static_folder='static')


load_dotenv()

app = Flask(__name__)
CORS(app)

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
sentence_model = SentenceTransformer('all-MiniLM-L6-v2')
document_text = ""

def extract_text_from_pdf(pdf_content):
    with fitz.open(stream=pdf_content, filetype="pdf") as doc:
        text = ""
        for page in doc:
            text += page.get_text()
    return text
def extract_text_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    txt = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt += shape.text + "\n"
    return txt
def extract_text_from_docx(docx_file):
    """
    Extract text from a .docx file.
    
    Args:
    - docx_file: Path to the .docx file
    
    Returns:
    - text: Extracted text from the .docx file
    """
    try:
        doc = docx.Document(docx_file)
        text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        print(f"Error extracting text from .docx file: {e}")
        return None
    
def extract_text_from_tex(tex_file):
    try:
        text = textract.process(tex_file).decode("utf-8")
        return text
    except Exception as e:
        print(f"Error extracting text from .tex file: {e}")
        print(type(text))
        return None


def extract_text_from_doc(doc_file):
    """
    Extract text from a .doc file.
    
    Args:
    - doc_file: Path to the .doc file
    
    Returns:
    - text: Extracted text from the .doc file
    """
    try:
        # Use textract to extract text from .doc files
        text = textract.process(doc_file).decode("utf-8")
        return text
    except Exception as e:
        print(f"Error extracting text from .doc file: {e}")
        return None


@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'document' not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files['document']    
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    if file:
        file_extension = file.filename.rsplit('.', 1)[1].lower()
        
        global document_text
        if file_extension == 'pdf':
            document_text = extract_text_from_pdf(file.read())
            return jsonify({"text": document_text}), 200
        elif file_extension == 'docx':
            document_text = extract_text_from_docx(file)
            return jsonify({"text": document_text}), 200
        elif file_extension == 'tex':
            document_text = extract_text_from_tex(file)
            return jsonify({"text": document_text}), 200
        elif file_extension == 'doc':
            document_text = extract_text_from_doc(file)
            return jsonify({"text": document_text}), 200
        elif file_extension in ['ppt', 'pptx']:
            document_text = extract_text_from_ppt(file)  
            return jsonify({"text": document_text}), 200
        else:
            return jsonify({"error": "Unsupported file type"}), 400
    else:
        return jsonify({"error": "Unsupported file type"}), 400

@app.route('/ask', methods=['POST'])
def ask_question():
    data = request.get_json()
    question =  data.get("question") + "answer in context of above research paper provided in 10 lines also state apporopriate lines from the research paper.If query is not in the context of research paper then respond ask user to enter appropriate query "
    
    global document_text
    # Check if document text is available
    if not document_text:
        return jsonify({"answer": "Please upload a document first."}), 400
    
    # Use ChatGPT API to generate response
    response = requests.post(
        "https://api.openai.com/v1/chat/completions",
        headers={
            "Authorization": f"Bearer {OPENAI_API_KEY}",
            "Content-Type": "application/json",
        },
        json={
            "model": "gpt-3.5-turbo",
            "messages": [
                {"role": "user", "content": document_text},
                {"role": "assistant", "content": question},
            ],
        }
    )
    print(response.text)
    if response.status_code == 200:
        # answer = response.json().get("choices")[0]['message']['content'].get("text", "No answer found.").strip()
        answer = response.json().get("choices")[0]['message']['content']
    else:
        answer = "Failed to get a response from OpenAI API."

    return jsonify({"answer": answer}), 200

if __name__ == '_main_':
    app.run(debug=True)